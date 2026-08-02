"""从三套资料提取带定位和哈希的公式证据。"""

import argparse
import hashlib
import json
import sys
from dataclasses import dataclass
from pathlib import Path

import fitz
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation

SKILL_ROOT = Path(__file__).parents[1]
sys.path.insert(0, str(SKILL_ROOT / "src"))

from cashflow_main.input_adapter import relevant_sheet_names

SOURCE_FILES = {
    "first_workbook": "案例和底稿1/00-现金流案例.xlsx",
    "first_document": "案例和底稿1/现金流量表4.0.docx",
    "second_slides": "案例和底稿2/现金流量表编制与复核技巧-谢海林.pptx",
    "second_workbook": "案例和底稿2/现金流量表编制与复核技巧-谢海林.xlsx",
    "third_workbook": "案例和底稿3/05 I-现金流量表审计底稿.xlsx",
    "knowledge_base": "知识库_现金流量表相关内容汇编.md",
    "pdf_article": "chapter23.pdf",
}


@dataclass(frozen=True)
class ExtractionResult:
    evidence: tuple[dict[str, object], ...]
    hidden_sheets_read: tuple[str, ...] = ()


def _evidence(source_id: str, locator: str, content: str) -> dict[str, object]:
    return {
        "source_id": source_id,
        "locator": locator,
        "content": content,
        "content_hash": hashlib.sha256(content.encode("utf-8")).hexdigest(),
    }


def extract_xlsx(path: Path, source_id: str) -> ExtractionResult:
    formulas = load_workbook(path, data_only=False, read_only=False)
    values = load_workbook(path, data_only=True, read_only=False)
    sheet_names = relevant_sheet_names(path)
    hidden = tuple(
        name for name in sheet_names if formulas[name].sheet_state != "visible"
    )
    records = []
    for sheet_name in sheet_names:
        formula_sheet = formulas[sheet_name]
        value_sheet = values[sheet_name]
        for row in formula_sheet.iter_rows():
            for cell in row:
                if cell.value in (None, ""):
                    continue
                displayed = value_sheet[cell.coordinate].value
                content = (
                    f"公式={cell.value}|显示值={displayed}"
                    if isinstance(cell.value, str) and cell.value.startswith("=")
                    else str(cell.value)
                )
                records.append(
                    _evidence(
                        source_id,
                        f"{path.name}:{sheet_name}!{cell.coordinate}",
                        content,
                    )
                )
    formulas.close()
    values.close()
    return ExtractionResult(tuple(records), hidden)


def extract_docx(path: Path, source_id: str) -> ExtractionResult:
    document = Document(path)
    records = []
    for index, paragraph in enumerate(document.paragraphs, start=1):
        text = paragraph.text.strip()
        if text:
            records.append(_evidence(source_id, f"{path.name}:段落{index}", text))
    for table_index, table in enumerate(document.tables, start=1):
        for row_index, row in enumerate(table.rows, start=1):
            text = " | ".join(cell.text.strip() for cell in row.cells).strip(" |")
            if text:
                records.append(
                    _evidence(
                        source_id,
                        f"{path.name}:表{table_index}-行{row_index}",
                        text,
                    )
                )
    return ExtractionResult(tuple(records))


def extract_pptx(path: Path, source_id: str) -> ExtractionResult:
    presentation = Presentation(path)
    records = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        for shape_index, shape in enumerate(slide.shapes, start=1):
            if not getattr(shape, "has_text_frame", False):
                continue
            text = shape.text.strip()
            if text:
                records.append(
                    _evidence(
                        source_id,
                        f"{path.name}:页{slide_index}-文本框{shape_index}",
                        text,
                    )
                )
    return ExtractionResult(tuple(records))


def extract_pdf(path: Path, source_id: str) -> ExtractionResult:
    document = fitz.open(path)
    records = []
    for page_index, page in enumerate(document, start=1):
        text = page.get_text("text").strip()
        if text:
            records.append(
                _evidence(source_id, f"{path.name}:页{page_index}", text)
            )
    document.close()
    return ExtractionResult(tuple(records))


def extract_markdown(path: Path, source_id: str) -> ExtractionResult:
    records = []
    for line_number, line in enumerate(
        path.read_text(encoding="utf-8-sig").splitlines(),
        start=1,
    ):
        text = line.strip()
        if text:
            records.append(
                _evidence(source_id, f"{path.name}:行{line_number}", text)
            )
    return ExtractionResult(tuple(records))


def extract_file(path: Path, source_id: str) -> ExtractionResult:
    extractors = {
        ".xlsx": extract_xlsx,
        ".docx": extract_docx,
        ".pptx": extract_pptx,
        ".pdf": extract_pdf,
        ".md": extract_markdown,
    }
    try:
        extractor = extractors[path.suffix.lower()]
    except KeyError as exc:
        raise ValueError(f"不支持的证据文件：{path}") from exc
    return extractor(path, source_id)


def extract_sources(source_root: Path, output_jsonl: Path) -> dict[str, object]:
    records = []
    source_summary = {}
    for source_id, relative in SOURCE_FILES.items():
        path = source_root / Path(relative)
        if not path.is_file():
            raise FileNotFoundError(f"源资料不存在：{path}")
        result = extract_file(path, source_id)
        records.extend(result.evidence)
        source_summary[source_id] = {
            "path": str(path),
            "sha256": hashlib.sha256(path.read_bytes()).hexdigest(),
            "size": path.stat().st_size,
            "evidence_count": len(result.evidence),
            "hidden_sheets_read": result.hidden_sheets_read,
        }
    output_jsonl.parent.mkdir(parents=True, exist_ok=True)
    with output_jsonl.open("w", encoding="utf-8-sig", newline="\n") as handle:
        for record in records:
            handle.write(json.dumps(record, ensure_ascii=False) + "\n")
    return source_summary


def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("--source-root", type=Path, required=True)
    parser.add_argument("--output", type=Path, required=True)
    parser.add_argument("--summary", type=Path)
    args = parser.parse_args()
    summary = extract_sources(args.source_root, args.output)
    if args.summary:
        args.summary.write_text(
            json.dumps(summary, ensure_ascii=False, indent=2),
            encoding="utf-8-sig",
        )
    print(
        f"已提取{sum(item['evidence_count'] for item in summary.values())}条证据"
    )
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
